#!/usr/bin/env python3
"""
Presentation Generator - Multi-format slide deck generation from Markdown.

Generates PDF, PPTX, and HTML (Reveal.js) presentations from markdown source files.
Supports multiple design styles and optional research content integration.

Usage:
    python generate.py --input source.md --output ./output --formats pdf,html,pptx --style systematic-velocity
"""

import argparse
import json
import os
import re
import sys
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Optional

import yaml

# Design style configurations
STYLES = {
    "systematic-velocity": {
        "name": "Systematic Velocity",
        "background": "#1A1A1A",
        "accent": "#DA7756",
        "text": "#EBEBEB",
        "text_secondary": "#A0A0A0",
        "heading_font": "Inter",
        "body_font": "Inter",
        "code_font": "JetBrains Mono",
    },
    "clean-minimal": {
        "name": "Clean Minimal",
        "background": "#FFFFFF",
        "accent": "#2563EB",
        "text": "#1F2937",
        "text_secondary": "#6B7280",
        "heading_font": "Inter",
        "body_font": "Inter",
        "code_font": "JetBrains Mono",
    },
    "technical-blueprint": {
        "name": "Technical Blueprint",
        "background": "#0F172A",
        "accent": "#06B6D4",
        "text": "#E2E8F0",
        "text_secondary": "#94A3B8",
        "heading_font": "Fira Code",
        "body_font": "Fira Code",
        "code_font": "Fira Code",
    },
    "marketing-bold": {
        "name": "Marketing Bold",
        "background": "#FFFFFF",
        "accent": "#8B5CF6",
        "accent_gradient": "linear-gradient(135deg, #8B5CF6, #EC4899)",
        "text": "#111827",
        "text_secondary": "#6B7280",
        "heading_font": "Poppins",
        "body_font": "Inter",
        "code_font": "JetBrains Mono",
    },
}


@dataclass
class SlideContent:
    """Represents a single slide."""

    title: str = ""
    content: list = field(default_factory=list)
    notes: str = ""
    slide_type: str = "content"  # title, content, section, metrics, code, image
    metrics: list = field(default_factory=list)
    columns: list = field(default_factory=list)
    code_blocks: list = field(default_factory=list)
    images: list = field(default_factory=list)
    cta: list = field(default_factory=list)


@dataclass
class Presentation:
    """Represents a complete presentation."""

    title: str = ""
    subtitle: str = ""
    author: str = ""
    date: str = ""
    style: str = "systematic-velocity"
    formats: list = field(default_factory=list)
    slides: list = field(default_factory=list)
    research_content: str = ""


def parse_frontmatter(content: str) -> tuple[dict, str]:
    """Extract YAML frontmatter from markdown content."""
    if content.startswith("---"):
        parts = content.split("---", 2)
        if len(parts) >= 3:
            try:
                frontmatter = yaml.safe_load(parts[1])
                return frontmatter or {}, parts[2].strip()
            except yaml.YAMLError:
                pass
    return {}, content


def parse_slide_content(slide_text: str) -> SlideContent:
    """Parse a single slide's markdown content."""
    slide = SlideContent()
    lines = slide_text.strip().split("\n")

    current_section = "content"
    content_lines = []
    notes_lines = []
    in_code_block = False
    code_block_lines = []
    code_lang = ""

    for line in lines:
        # Handle code blocks
        if line.startswith("```"):
            if in_code_block:
                slide.code_blocks.append({"lang": code_lang, "code": "\n".join(code_block_lines)})
                code_block_lines = []
                in_code_block = False
            else:
                in_code_block = True
                code_lang = line[3:].strip() or "text"
            continue

        if in_code_block:
            code_block_lines.append(line)
            continue

        # Handle notes
        if line.strip() == "::: notes":
            current_section = "notes"
            continue
        elif line.strip() == ":::":
            current_section = "content"
            continue

        # Handle columns
        if line.strip() == "::: columns":
            current_section = "columns"
            slide.columns = []
            continue
        elif line.strip().startswith(":::: column"):
            slide.columns.append([])
            continue
        elif line.strip() == "::::":
            continue

        # Handle metrics
        metric_match = re.match(r"<!--\s*metric:\s*value=\"([^\"]+)\"\s*label=\"([^\"]+)\"", line)
        if metric_match:
            slide.metrics.append({"value": metric_match.group(1), "label": metric_match.group(2)})
            continue

        # Handle CTAs
        cta_match = re.match(r"<!--\s*cta:\s*text=\"([^\"]+)\"\s*url=\"([^\"]+)\"", line)
        if cta_match:
            slide.cta.append({"text": cta_match.group(1), "url": cta_match.group(2)})
            continue

        # Handle images
        img_match = re.match(r"!\[([^\]]*)\]\(([^)]+)\)", line)
        if img_match:
            slide.images.append({"alt": img_match.group(1), "src": img_match.group(2)})
            continue

        # Handle headings
        if line.startswith("# "):
            slide.title = line[2:].strip()
            slide.slide_type = "section"
            continue
        elif line.startswith("## "):
            slide.title = line[3:].strip()
            continue

        # Add to appropriate section
        if current_section == "notes":
            notes_lines.append(line)
        elif current_section == "columns" and slide.columns:
            slide.columns[-1].append(line)
        else:
            content_lines.append(line)

    slide.content = [l for l in content_lines if l.strip()]
    slide.notes = "\n".join(notes_lines).strip()

    # Determine slide type
    if slide.metrics:
        slide.slide_type = "metrics"
    elif slide.code_blocks:
        slide.slide_type = "code"
    elif slide.images:
        slide.slide_type = "image"
    elif not slide.title and not slide.content:
        slide.slide_type = "blank"

    return slide


def parse_markdown(content: str) -> Presentation:
    """Parse markdown content into a Presentation object."""
    frontmatter, body = parse_frontmatter(content)

    pres = Presentation(
        title=frontmatter.get("title", "Untitled Presentation"),
        subtitle=frontmatter.get("subtitle", ""),
        author=frontmatter.get("author", ""),
        date=frontmatter.get("date", datetime.now().strftime("%Y-%m-%d")),
        style=frontmatter.get("style", "systematic-velocity"),
        formats=frontmatter.get("format", ["pdf"]),
    )

    # Split by slide separator
    slide_texts = re.split(r"\n---\n", body)

    for i, slide_text in enumerate(slide_texts):
        if not slide_text.strip():
            continue

        slide = parse_slide_content(slide_text)

        # First slide is title slide
        if i == 0 and not slide.title:
            slide.title = pres.title
            slide.slide_type = "title"
            slide.content = [pres.subtitle] if pres.subtitle else []

        pres.slides.append(slide)

    return pres


def load_research_content(research_dir: str) -> str:
    """Load research content from directory."""
    research_path = Path(research_dir)
    if not research_path.exists():
        return ""

    content = []
    for md_file in research_path.glob("*.md"):
        content.append(md_file.read_text())

    return "\n\n".join(content)


def generate_pdf(presentation: Presentation, output_path: str, style: dict) -> None:
    """Generate PDF presentation using ReportLab."""
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import LETTER, landscape
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer
    except ImportError:
        print("ReportLab not installed. Skipping PDF generation.")
        print("Install with: pip install reportlab")
        return

    page_width, page_height = landscape(LETTER)
    doc = SimpleDocTemplate(
        output_path,
        pagesize=landscape(LETTER),
        leftMargin=0.75 * inch,
        rightMargin=0.75 * inch,
        topMargin=0.5 * inch,
        bottomMargin=0.5 * inch,
    )

    # Convert hex colors
    def hex_to_color(hex_color: str) -> colors.Color:
        hex_color = hex_color.lstrip("#")
        r, g, b = tuple(int(hex_color[i : i + 2], 16) / 255 for i in (0, 2, 4))
        return colors.Color(r, g, b)

    bg_color = hex_to_color(style["background"])
    text_color = hex_to_color(style["text"])
    accent_color = hex_to_color(style["accent"])

    # Create styles
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "SlideTitle",
        parent=styles["Heading1"],
        fontSize=36,
        textColor=text_color,
        spaceAfter=30,
        alignment=1,  # Center
    )

    heading_style = ParagraphStyle(
        "SlideHeading",
        parent=styles["Heading2"],
        fontSize=28,
        textColor=accent_color,
        spaceAfter=20,
    )

    body_style = ParagraphStyle(
        "SlideBody",
        parent=styles["Normal"],
        fontSize=18,
        textColor=text_color,
        spaceAfter=12,
        leading=24,
    )

    metric_value_style = ParagraphStyle(
        "MetricValue",
        parent=styles["Normal"],
        fontSize=48,
        textColor=accent_color,
        alignment=1,
        spaceAfter=8,
    )

    metric_label_style = ParagraphStyle(
        "MetricLabel",
        parent=styles["Normal"],
        fontSize=16,
        textColor=hex_to_color(style["text_secondary"]),
        alignment=1,
        spaceAfter=20,
    )

    story = []

    for i, slide in enumerate(presentation.slides):
        if i > 0:
            story.append(Spacer(1, page_height - 2 * inch))  # Page break

        # Title
        if slide.title:
            if slide.slide_type == "title":
                story.append(Spacer(1, 2 * inch))
                story.append(Paragraph(slide.title, title_style))
                if presentation.subtitle:
                    story.append(Paragraph(presentation.subtitle, body_style))
                story.append(Spacer(1, inch))
                if presentation.author:
                    story.append(Paragraph(f"{presentation.author} | {presentation.date}", body_style))
            else:
                story.append(Paragraph(slide.title, heading_style))

        # Metrics
        for metric in slide.metrics:
            story.append(Paragraph(metric["value"], metric_value_style))
            story.append(Paragraph(metric["label"], metric_label_style))

        # Content
        for line in slide.content:
            if line.startswith("- "):
                story.append(Paragraph(f"• {line[2:]}", body_style))
            elif line.startswith("> "):
                quote_style = ParagraphStyle(
                    "Quote",
                    parent=body_style,
                    fontSize=20,
                    textColor=accent_color,
                    leftIndent=30,
                    fontName="Times-Italic",
                )
                story.append(Paragraph(line[2:], quote_style))
            elif line.strip():
                story.append(Paragraph(line, body_style))

        # Code blocks
        for code_block in slide.code_blocks:
            code_style = ParagraphStyle(
                "Code",
                parent=styles["Code"],
                fontSize=12,
                textColor=text_color,
                backColor=hex_to_color("#2D2D2D"),
                leftIndent=20,
                rightIndent=20,
                spaceBefore=10,
                spaceAfter=10,
            )
            code_text = code_block["code"].replace("\n", "<br/>")
            story.append(Paragraph(f"<pre>{code_text}</pre>", code_style))

    doc.build(story)
    print(f"Generated PDF: {output_path}")


def generate_pptx(presentation: Presentation, output_path: str, style: dict) -> None:
    """Generate PowerPoint presentation using python-pptx."""
    try:
        from pptx import Presentation as PPTXPresentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError:
        print("python-pptx not installed. Skipping PPTX generation.")
        print("Install with: pip install python-pptx")
        return

    def hex_to_rgb(hex_color: str) -> RGBColor:
        hex_color = hex_color.lstrip("#")
        return RGBColor(*[int(hex_color[i : i + 2], 16) for i in (0, 2, 4)])

    prs = PPTXPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_color = hex_to_rgb(style["background"])
    text_color = hex_to_rgb(style["text"])
    accent_color = hex_to_rgb(style["accent"])

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for slide_data in presentation.slides:
        pptx_slide = prs.slides.add_slide(blank_layout)

        # Set background
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        y_pos = Inches(0.75)

        # Title
        if slide_data.title:
            title_box = pptx_slide.shapes.add_textbox(Inches(0.75), y_pos, Inches(11.833), Inches(1))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = slide_data.title
            title_para.font.size = Pt(44 if slide_data.slide_type == "title" else 36)
            title_para.font.color.rgb = accent_color if slide_data.slide_type != "title" else text_color
            title_para.font.bold = True
            y_pos += Inches(1.2)

        # Subtitle for title slide
        if slide_data.slide_type == "title" and presentation.subtitle:
            sub_box = pptx_slide.shapes.add_textbox(Inches(0.75), y_pos, Inches(11.833), Inches(0.75))
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = presentation.subtitle
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = hex_to_rgb(style["text_secondary"])
            y_pos += Inches(1.5)

            if presentation.author:
                auth_box = pptx_slide.shapes.add_textbox(Inches(0.75), Inches(6), Inches(11.833), Inches(0.5))
                auth_frame = auth_box.text_frame
                auth_para = auth_frame.paragraphs[0]
                auth_para.text = f"{presentation.author} | {presentation.date}"
                auth_para.font.size = Pt(16)
                auth_para.font.color.rgb = hex_to_rgb(style["text_secondary"])

        # Metrics
        if slide_data.metrics:
            metric_width = Inches(11.833 / len(slide_data.metrics))
            for i, metric in enumerate(slide_data.metrics):
                x_pos = Inches(0.75) + (i * metric_width)

                # Value
                val_box = pptx_slide.shapes.add_textbox(x_pos, y_pos, metric_width, Inches(1))
                val_frame = val_box.text_frame
                val_para = val_frame.paragraphs[0]
                val_para.text = metric["value"]
                val_para.font.size = Pt(56)
                val_para.font.color.rgb = accent_color
                val_para.font.bold = True

                # Label
                lab_box = pptx_slide.shapes.add_textbox(x_pos, y_pos + Inches(1.2), metric_width, Inches(0.5))
                lab_frame = lab_box.text_frame
                lab_para = lab_frame.paragraphs[0]
                lab_para.text = metric["label"]
                lab_para.font.size = Pt(18)
                lab_para.font.color.rgb = hex_to_rgb(style["text_secondary"])

            y_pos += Inches(2)

        # Content
        if slide_data.content:
            content_box = pptx_slide.shapes.add_textbox(Inches(0.75), y_pos, Inches(11.833), Inches(4))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for j, line in enumerate(slide_data.content):
                if j == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                if line.startswith("- "):
                    para.text = f"• {line[2:]}"
                    para.level = 0
                elif line.startswith("> "):
                    para.text = line[2:]
                    para.font.italic = True
                else:
                    para.text = line

                para.font.size = Pt(20)
                para.font.color.rgb = text_color
                para.space_after = Pt(12)

    prs.save(output_path)
    print(f"Generated PPTX: {output_path}")


def generate_html(presentation: Presentation, output_dir: str, style: dict) -> None:
    """Generate HTML presentation using Reveal.js."""
    try:
        from jinja2 import Template
    except ImportError:
        print("Jinja2 not installed. Skipping HTML generation.")
        print("Install with: pip install jinja2")
        return

    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)

    html_template = Template(
        """<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{{ title }}</title>
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/reveal.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/theme/black.css">
    <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/highlight.js@11.9.0/styles/github-dark.min.css">
    <style>
        :root {
            --r-background-color: {{ style.background }};
            --r-main-color: {{ style.text }};
            --r-heading-color: {{ style.accent }};
            --r-link-color: {{ style.accent }};
        }
        .reveal {
            font-family: '{{ style.body_font }}', sans-serif;
        }
        .reveal h1, .reveal h2, .reveal h3 {
            font-family: '{{ style.heading_font }}', sans-serif;
            color: {{ style.accent }};
        }
        .reveal .slides section {
            text-align: left;
        }
        .reveal .title-slide {
            text-align: center;
        }
        .reveal .title-slide h1 {
            color: {{ style.text }};
            font-size: 2.5em;
        }
        .reveal .metric {
            text-align: center;
            display: inline-block;
            margin: 0 2em;
        }
        .reveal .metric-value {
            font-size: 3em;
            font-weight: bold;
            color: {{ style.accent }};
            font-family: '{{ style.code_font }}', monospace;
        }
        .reveal .metric-label {
            font-size: 0.9em;
            color: {{ style.text_secondary }};
        }
        .reveal pre code {
            font-family: '{{ style.code_font }}', monospace;
            font-size: 0.8em;
            padding: 1em;
            border-radius: 8px;
        }
        .reveal blockquote {
            border-left: 4px solid {{ style.accent }};
            padding-left: 1em;
            font-style: italic;
        }
        .reveal .columns {
            display: flex;
            gap: 2em;
        }
        .reveal .column {
            flex: 1;
        }
        .reveal .cta-button {
            display: inline-block;
            padding: 0.75em 2em;
            background: {{ style.accent }};
            color: {{ style.background }};
            text-decoration: none;
            border-radius: 8px;
            font-weight: bold;
            margin: 0.5em;
        }
        .reveal .speaker-notes {
            display: none;
        }
    </style>
</head>
<body>
    <div class="reveal">
        <div class="slides">
            {% for slide in slides %}
            <section class="{{ 'title-slide' if slide.slide_type == 'title' else '' }}">
                {% if slide.title %}
                <h{{ '1' if slide.slide_type in ['title', 'section'] else '2' }}>{{ slide.title }}</h{{ '1' if slide.slide_type in ['title', 'section'] else '2' }}>
                {% endif %}

                {% if slide.slide_type == 'title' and subtitle %}
                <p style="color: {{ style.text_secondary }}">{{ subtitle }}</p>
                <p style="font-size: 0.7em; color: {{ style.text_secondary }}; margin-top: 2em;">{{ author }} | {{ date }}</p>
                {% endif %}

                {% if slide.metrics %}
                <div style="text-align: center; margin-top: 2em;">
                {% for metric in slide.metrics %}
                    <div class="metric">
                        <div class="metric-value">{{ metric.value }}</div>
                        <div class="metric-label">{{ metric.label }}</div>
                    </div>
                {% endfor %}
                </div>
                {% endif %}

                {% if slide.content %}
                <ul>
                {% for line in slide.content %}
                    {% if line.startswith('- ') %}
                    <li>{{ line[2:] }}</li>
                    {% elif line.startswith('> ') %}
                    <blockquote>{{ line[2:] }}</blockquote>
                    {% elif line.strip() %}
                    <p>{{ line }}</p>
                    {% endif %}
                {% endfor %}
                </ul>
                {% endif %}

                {% if slide.code_blocks %}
                {% for code in slide.code_blocks %}
                <pre><code class="language-{{ code.lang }}">{{ code.code }}</code></pre>
                {% endfor %}
                {% endif %}

                {% if slide.columns %}
                <div class="columns">
                {% for col in slide.columns %}
                    <div class="column">
                    {% for line in col %}
                        {% if line.startswith('### ') %}
                        <h3>{{ line[4:] }}</h3>
                        {% elif line.startswith('- ') %}
                        <p>• {{ line[2:] }}</p>
                        {% elif line.strip() %}
                        <p>{{ line }}</p>
                        {% endif %}
                    {% endfor %}
                    </div>
                {% endfor %}
                </div>
                {% endif %}

                {% if slide.cta %}
                <div style="text-align: center; margin-top: 2em;">
                {% for cta in slide.cta %}
                    <a href="{{ cta.url }}" class="cta-button">{{ cta.text }}</a>
                {% endfor %}
                </div>
                {% endif %}

                {% if slide.notes %}
                <aside class="notes">{{ slide.notes }}</aside>
                {% endif %}
            </section>
            {% endfor %}
        </div>
    </div>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/dist/reveal.js"></script>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/plugin/notes/notes.js"></script>
    <script src="https://cdn.jsdelivr.net/npm/reveal.js@4.6.0/plugin/highlight/highlight.js"></script>
    <script>
        Reveal.initialize({
            hash: true,
            plugins: [RevealNotes, RevealHighlight],
            transition: 'slide',
            backgroundTransition: 'fade'
        });
    </script>
</body>
</html>
"""
    )

    slides_data = []
    for slide in presentation.slides:
        slides_data.append(
            {
                "title": slide.title,
                "content": slide.content,
                "notes": slide.notes,
                "slide_type": slide.slide_type,
                "metrics": slide.metrics,
                "columns": slide.columns,
                "code_blocks": slide.code_blocks,
                "images": slide.images,
                "cta": slide.cta,
            }
        )

    html_content = html_template.render(
        title=presentation.title,
        subtitle=presentation.subtitle,
        author=presentation.author,
        date=presentation.date,
        style=style,
        slides=slides_data,
    )

    index_path = output_path / "index.html"
    index_path.write_text(html_content)
    print(f"Generated HTML: {index_path}")


def generate_markdown(presentation: Presentation, output_path: str) -> None:
    """Generate clean markdown output (useful for editing)."""
    lines = []

    # Frontmatter
    lines.append("---")
    lines.append(f'title: "{presentation.title}"')
    if presentation.subtitle:
        lines.append(f'subtitle: "{presentation.subtitle}"')
    if presentation.author:
        lines.append(f'author: "{presentation.author}"')
    lines.append(f"date: {presentation.date}")
    lines.append(f"style: {presentation.style}")
    lines.append("format:")
    for fmt in presentation.formats:
        lines.append(f"  - {fmt}")
    lines.append("---")
    lines.append("")

    for i, slide in enumerate(presentation.slides):
        if i > 0:
            lines.append("")
            lines.append("---")
            lines.append("")

        if slide.title:
            prefix = "#" if slide.slide_type in ["title", "section"] else "##"
            lines.append(f"{prefix} {slide.title}")
            lines.append("")

        for metric in slide.metrics:
            lines.append(f'<!-- metric: value="{metric["value"]}" label="{metric["label"]}" -->')

        for line in slide.content:
            lines.append(line)

        for code in slide.code_blocks:
            lines.append(f"```{code['lang']}")
            lines.append(code["code"])
            lines.append("```")

        if slide.columns:
            lines.append("::: columns")
            for col in slide.columns:
                lines.append(":::: column")
                for line in col:
                    lines.append(line)
                lines.append("::::")
            lines.append(":::")

        for cta in slide.cta:
            lines.append(f'<!-- cta: text="{cta["text"]}" url="{cta["url"]}" -->')

        if slide.notes:
            lines.append("")
            lines.append("::: notes")
            lines.append(slide.notes)
            lines.append(":::")

    Path(output_path).write_text("\n".join(lines))
    print(f"Generated Markdown: {output_path}")


def main():
    parser = argparse.ArgumentParser(description="Generate presentations from markdown")
    parser.add_argument("--input", "-i", required=True, help="Input markdown file")
    parser.add_argument("--output", "-o", required=True, help="Output directory or file")
    parser.add_argument(
        "--formats",
        "-f",
        default="pdf",
        help="Comma-separated formats: pdf,html,pptx,md,all",
    )
    parser.add_argument(
        "--style",
        "-s",
        default="systematic-velocity",
        choices=list(STYLES.keys()),
        help="Design style",
    )
    parser.add_argument(
        "--research-dir",
        "-r",
        default="",
        help="Directory containing research content",
    )

    args = parser.parse_args()

    # Load input
    input_path = Path(args.input)
    if not input_path.exists():
        print(f"Error: Input file not found: {args.input}")
        sys.exit(1)

    content = input_path.read_text()

    # Load research if provided
    research_content = ""
    if args.research_dir:
        research_content = load_research_content(args.research_dir)

    # Parse presentation
    presentation = parse_markdown(content)
    presentation.research_content = research_content

    # Override style if specified
    if args.style:
        presentation.style = args.style

    style = STYLES.get(presentation.style, STYLES["systematic-velocity"])

    # Parse formats
    formats = args.formats.lower().split(",")
    if "all" in formats:
        formats = ["pdf", "html", "pptx", "md"]

    # Setup output
    output_path = Path(args.output)
    if output_path.suffix:
        # Single file output
        output_dir = output_path.parent
        base_name = output_path.stem
    else:
        # Directory output
        output_dir = output_path
        base_name = input_path.stem

    output_dir.mkdir(parents=True, exist_ok=True)

    # Generate outputs
    for fmt in formats:
        fmt = fmt.strip()
        if fmt == "pdf":
            generate_pdf(presentation, str(output_dir / f"{base_name}.pdf"), style)
        elif fmt == "pptx":
            generate_pptx(presentation, str(output_dir / f"{base_name}.pptx"), style)
        elif fmt == "html":
            generate_html(presentation, str(output_dir / base_name), style)
        elif fmt == "md":
            generate_markdown(presentation, str(output_dir / f"{base_name}.md"))

    print(f"\nGeneration complete! Output in: {output_dir}")


if __name__ == "__main__":
    main()
